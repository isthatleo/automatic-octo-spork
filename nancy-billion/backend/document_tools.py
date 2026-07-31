"""Real Office document / PDF creation -- the write-side counterpart to
doc_extraction.py's read-only extraction. Confirmed live (2026-07-31
Hermes-parity audit): Nancy could read Word/PDF attachments but had no way
to produce one, unlike Hermes' "Office docs (Word, Excel, PowerPoint), PDF
manipulation" skill. Word/Excel/PDF/PowerPoint all covered here.

Every tool writes through file_access.safe_path -- the same allowlist
sandbox every other file-writing tool in this codebase goes through, since
these are just another way to put bytes on disk.
"""
from __future__ import annotations

import logging
from typing import Any, Dict, List, Optional

from file_access import PathNotAllowed, safe_path

logger = logging.getLogger(__name__)

DOCUMENT_TOOLS = [
    {
        "name": "create_word_document",
        "description": (
            "Create a real .docx Word document with a title and a series of paragraphs/headings. "
            "Writes to the given path (created if it doesn't exist)."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Output file path, e.g. 'reports/summary.docx'."},
                "title": {"type": "string"},
                "sections": {
                    "type": "array",
                    "description": "Each item is a paragraph of body text, or {\"heading\": \"...\"} for a section heading.",
                    "items": {},
                },
            },
            "required": ["path", "title"],
        },
    },
    {
        "name": "create_excel_spreadsheet",
        "description": (
            "Create a real .xlsx spreadsheet from tabular data -- one sheet, first row as headers. "
            "Writes to the given path (created if it doesn't exist)."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Output file path, e.g. 'reports/data.xlsx'."},
                "sheet_name": {"type": "string"},
                "headers": {"type": "array", "items": {"type": "string"}},
                "rows": {
                    "type": "array",
                    "description": "Each item is a row: an array of cell values (string/number).",
                    "items": {"type": "array"},
                },
            },
            "required": ["path", "headers", "rows"],
        },
    },
    {
        "name": "create_pdf_document",
        "description": (
            "Create a real, printable .pdf with a title and a series of paragraphs/headings. "
            "Writes to the given path (created if it doesn't exist)."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Output file path, e.g. 'reports/summary.pdf'."},
                "title": {"type": "string"},
                "sections": {
                    "type": "array",
                    "description": "Each item is a paragraph of body text, or {\"heading\": \"...\"} for a section heading.",
                    "items": {},
                },
            },
            "required": ["path", "title"],
        },
    },
    {
        "name": "create_powerpoint_presentation",
        "description": (
            "Create a real .pptx PowerPoint presentation -- a title slide, then one slide per "
            "entry with a heading and bullet points. Writes to the given path (created if it "
            "doesn't exist)."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Output file path, e.g. 'reports/deck.pptx'."},
                "title": {"type": "string"},
                "subtitle": {"type": "string"},
                "slides": {
                    "type": "array",
                    "description": "Each item: {\"heading\": \"...\", \"bullets\": [\"...\", \"...\"]}.",
                    "items": {
                        "type": "object",
                        "properties": {
                            "heading": {"type": "string"},
                            "bullets": {"type": "array", "items": {"type": "string"}},
                        },
                    },
                },
            },
            "required": ["path", "title", "slides"],
        },
    },
]


def create_word_document(path: str, title: str, sections: Optional[List[Any]] = None) -> Dict[str, Any]:
    try:
        import docx
    except ImportError:
        return {"success": False, "error": "python-docx is not installed"}
    try:
        p = safe_path(path)
    except PathNotAllowed as e:
        return {"success": False, "error": str(e), "denied": True}

    document = docx.Document()
    document.add_heading(title, level=0)
    for item in (sections or []):
        if isinstance(item, dict) and "heading" in item:
            document.add_heading(str(item["heading"]), level=1)
        else:
            document.add_paragraph(str(item))

    p.parent.mkdir(parents=True, exist_ok=True)
    document.save(str(p))
    return {"success": True, "path": str(p)}


def create_excel_spreadsheet(
    path: str, headers: List[str], rows: List[List[Any]], sheet_name: str = "Sheet1"
) -> Dict[str, Any]:
    try:
        from openpyxl import Workbook
    except ImportError:
        return {"success": False, "error": "openpyxl is not installed"}
    try:
        p = safe_path(path)
    except PathNotAllowed as e:
        return {"success": False, "error": str(e), "denied": True}

    wb = Workbook()
    ws = wb.active
    ws.title = (sheet_name or "Sheet1")[:31]  # Excel's own 31-char sheet-name limit
    ws.append(list(headers))
    for row in rows:
        ws.append(list(row))

    p.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(p))
    return {"success": True, "path": str(p), "rows_written": len(rows)}


def create_pdf_document(path: str, title: str, sections: Optional[List[Any]] = None) -> Dict[str, Any]:
    try:
        from reportlab.lib.pagesizes import LETTER
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    except ImportError:
        return {"success": False, "error": "reportlab is not installed"}
    try:
        p = safe_path(path)
    except PathNotAllowed as e:
        return {"success": False, "error": str(e), "denied": True}

    styles = getSampleStyleSheet()
    story = [Paragraph(title, styles["Title"]), Spacer(1, 16)]
    for item in (sections or []):
        if isinstance(item, dict) and "heading" in item:
            story.append(Paragraph(str(item["heading"]), styles["Heading2"]))
        else:
            story.append(Paragraph(str(item), styles["BodyText"]))
        story.append(Spacer(1, 10))

    p.parent.mkdir(parents=True, exist_ok=True)
    SimpleDocTemplate(str(p), pagesize=LETTER).build(story)
    return {"success": True, "path": str(p)}


def create_powerpoint_presentation(
    path: str, title: str, slides: List[Dict[str, Any]], subtitle: str = ""
) -> Dict[str, Any]:
    try:
        from pptx import Presentation
    except ImportError:
        return {"success": False, "error": "python-pptx is not installed"}
    try:
        p = safe_path(path)
    except PathNotAllowed as e:
        return {"success": False, "error": str(e), "denied": True}

    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    if subtitle and len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = subtitle

    bullet_layout = prs.slide_layouts[1]
    for slide_spec in (slides or []):
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = str(slide_spec.get("heading", ""))
        body = slide.placeholders[1].text_frame
        bullets = slide_spec.get("bullets") or []
        for i, bullet in enumerate(bullets):
            para = body.paragraphs[0] if i == 0 else body.add_paragraph()
            para.text = str(bullet)

    p.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(p))
    return {"success": True, "path": str(p), "slide_count": len(slides or []) + 1}
